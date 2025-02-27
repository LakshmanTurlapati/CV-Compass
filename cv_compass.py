import json
import os
import hashlib
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import joblib
import tempfile
import shutil

# Configuration
RESUMES_DIR = "resumes"
CACHE_DIR = ".cache"
JD_DIR = "job_descriptions"

def load_resumes():
    """Load resumes from various file formats"""
    resumes = []
    filenames = []
    
    for filename in os.listdir(RESUMES_DIR):
        # Skip macOS system files
        if filename == ".DS_Store":
            continue
            
        path = os.path.join(RESUMES_DIR, filename)
        try:
            # PDF files
            if filename.lower().endswith('.pdf'):
                from PyPDF2 import PdfReader
                with open(path, 'rb') as f:
                    pdf = PdfReader(f)
                    text = "\n".join([page.extract_text() for page in pdf.pages])
            
            # Word documents
            elif filename.lower().endswith('.docx'):
                from docx import Document
                doc = Document(path)
                text = "\n".join([para.text for para in doc.paragraphs])
            
            # Plain text
            elif filename.lower().endswith('.txt'):
                with open(path, 'r', encoding='utf-8') as f:
                    text = f.read()
            
            # Excel files
            elif filename.lower().endswith(('.xls', '.xlsx')):
                import xlrd
                workbook = xlrd.open_workbook(path)
                text = ""
                for sheet in workbook.sheets():
                    for row in range(sheet.nrows):
                        text += " ".join(str(cell.value) for cell in sheet.row(row)) + "\n"
            
            # PowerPoint files
            elif filename.lower().endswith(('.ppt', '.pptx')):
                from pptx import Presentation
                prs = Presentation(path)
                text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            
            else:
                raise ValueError(f"Unsupported file format: {filename}")
            
            resumes.append(text)
            filenames.append(filename)
            
        except Exception as e:
            print(f"Error reading {filename}: {str(e)}")
            continue
            
    return filenames, resumes

def get_file_hash(filename):
    """Generate hash of file contents for cache invalidation"""
    with open(filename, 'rb') as f:
        return hashlib.md5(f.read()).hexdigest()

def generate_tfidf_matrix(resumes):
    """Create TF-IDF matrix and save vectorizer"""
    vectorizer = TfidfVectorizer(stop_words='english')
    tfidf_matrix = vectorizer.fit_transform(resumes)
    os.makedirs(CACHE_DIR, exist_ok=True)
    joblib.dump(vectorizer, os.path.join(CACHE_DIR, 'vectorizer.joblib'))
    joblib.dump(tfidf_matrix, os.path.join(CACHE_DIR, 'tfidf_matrix.joblib'))
    return tfidf_matrix, vectorizer

def load_or_generate_tfidf():
    """Load cached TF-IDF data or generate new"""
    cache_files = ['tfidf_matrix.joblib', 'vectorizer.joblib', 'resume_hashes.txt']
    if not all(os.path.exists(os.path.join(CACHE_DIR, f)) for f in cache_files):
        return generate_tfidf_matrix(*load_resumes()[1:])

    # Check if any resumes have changed
    current_hashes = {f: get_file_hash(os.path.join(RESUMES_DIR, f)) 
                     for f in os.listdir(RESUMES_DIR)}
    cached_hashes = joblib.load(os.path.join(CACHE_DIR, 'resume_hashes.txt'))
    
    if current_hashes != cached_hashes:
        return generate_tfidf_matrix(*load_resumes()[1:])
    
    return (joblib.load(os.path.join(CACHE_DIR, 'tfidf_matrix.joblib')),
            joblib.load(os.path.join(CACHE_DIR, 'vectorizer.joblib')))

def find_best_match(jd_file):
    """Main function to find best matching resume for a JD"""
    # Load JD
    with open(jd_file, 'r') as f:
        jd_data = json.load(f)
        jd_text = f"{jd_data['title']} {jd_data['description']} {jd_data['requirements']}"
    
    # Get TF-IDF data
    tfidf_matrix, vectorizer = load_or_generate_tfidf()
    
    # Transform JD
    jd_tfidf = vectorizer.transform([jd_text])
    
    # Calculate similarities
    similarities = cosine_similarity(jd_tfidf, tfidf_matrix)
    
    # Get best match
    filenames = load_resumes()[0]
    best_match_idx = similarities.argmax()
    
    # Print analysis
    print("\nTop Matching Resume Analysis:")
    print_important_terms(filenames[best_match_idx], jd_text, vectorizer)
    
    return filenames[best_match_idx]

def print_important_terms(resume_name, jd_text, vectorizer, num_terms=10):
    """Print key terms that influenced the match"""
    # Get resume text
    resume_idx = load_resumes()[0].index(resume_name)
    resume_text = load_resumes()[1][resume_idx]
    
    # Analyze important terms
    jd_vector = vectorizer.transform([jd_text])
    resume_vector = vectorizer.transform([resume_text])
    
    # Get overlapping terms with highest TF-IDF scores
    important_terms = {}
    for term in vectorizer.get_feature_names_out():
        jd_score = jd_vector[0, vectorizer.vocabulary_[term]]
        resume_score = resume_vector[0, vectorizer.vocabulary_[term]]
        if jd_score > 0 and resume_score > 0:
            important_terms[term] = (jd_score + resume_score) / 2
    
    # Sort and print top terms
    sorted_terms = sorted(important_terms.items(), key=lambda x: x[1], reverse=True)[:num_terms]
    
    # print(f"\nKey matching terms in '{resume_name}':")
    # for term, score in sorted_terms:
    #     print(f"- {term}: {score:.4f}")

# Example usage
if __name__ == "__main__":
    # jd_file = os.path.join(JD_DIR, "data_scientist.json")
    jd_file = os.path.join(JD_DIR, "full_stack_dev.json")

    print(f"Best match: {find_best_match(jd_file)}")
